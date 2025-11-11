"""Utilities for storing uploaded files and extracting readable context."""

from __future__ import annotations

import csv
import json
import mimetypes
import re
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List
from uuid import uuid4

from fastapi import HTTPException

from config.settings import settings

try:  # Optional dependency for PDF parsing
    from PyPDF2 import PdfReader  # type: ignore
except ImportError:  # pragma: no cover
    PdfReader = None

try:
    from docx import Document  # type: ignore
except ImportError:  # pragma: no cover
    Document = None

try:
    from openpyxl import load_workbook  # type: ignore
except ImportError:  # pragma: no cover
    load_workbook = None

try:
    from pptx import Presentation  # type: ignore
except ImportError:  # pragma: no cover
    Presentation = None

ATTACHMENTS_DIR = Path(settings.upload_dir) / "attachments"
ATTACHMENTS_DIR.mkdir(parents=True, exist_ok=True)

_SAFE_FILENAME_PATTERN = re.compile(r"[^A-Za-z0-9_.-]+")


def _sanitize_filename(filename: str) -> str:
    name = filename.strip().replace("/", "_").replace("\\", "_")
    sanitized = _SAFE_FILENAME_PATTERN.sub("_", name)
    return sanitized or "file"


def _generate_attachment_id() -> str:
    return uuid4().hex


def _ensure_allowed_extension(filename: str) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix not in settings.allowed_extensions:
        raise HTTPException(status_code=400, detail=f"不支持的文件类型: {suffix or 'unknown'}")
    return suffix


def _ensure_dependency(condition: bool, message: str) -> None:
    if not condition:
        raise HTTPException(status_code=500, detail=message)


def _read_text_file(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _extract_from_pdf(path: Path) -> str:
    _ensure_dependency(PdfReader is not None, "缺少 PyPDF2 依赖，无法解析 PDF 文件")
    reader = PdfReader(str(path))
    texts = []
    for page in reader.pages:
        page_text = page.extract_text() or ""
        texts.append(page_text)
    return "\n".join(texts)


def _extract_from_docx(path: Path) -> str:
    _ensure_dependency(Document is not None, "缺少 python-docx 依赖，无法解析 Word 文件")
    doc = Document(str(path))
    return "\n".join(paragraph.text for paragraph in doc.paragraphs)


def _extract_from_xlsx(path: Path) -> str:
    _ensure_dependency(load_workbook is not None, "缺少 openpyxl 依赖，无法解析 Excel 文件")
    workbook = load_workbook(filename=str(path), data_only=True, read_only=True)
    texts: List[str] = []
    row_limit = 50
    column_limit = 10
    for sheet in workbook.worksheets:
        texts.append(f"工作表: {sheet.title}")
        for row_index, row in enumerate(
            sheet.iter_rows(min_row=1, max_row=row_limit, values_only=True),
            start=1,
        ):
            cells = [str(cell) if cell is not None else "" for cell in row[:column_limit]]
            texts.append("\t".join(cells))
            if row_index >= row_limit:
                break
    return "\n".join(texts)


def _extract_from_pptx(path: Path) -> str:
    _ensure_dependency(Presentation is not None, "缺少 python-pptx 依赖，无法解析 PPT 文件")
    presentation = Presentation(str(path))
    texts: List[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts.append(f"幻灯片 {slide_index}")
        slide_text: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
        if slide_text:
            texts.append("\n".join(slide_text))
    return "\n".join(texts)


def _extract_from_csv(path: Path) -> str:
    with path.open("r", encoding="utf-8", errors="ignore", newline="") as csv_file:
        reader = csv.reader(csv_file)
        rows = []
        for index, row in enumerate(reader, start=1):
            rows.append(",".join(row))
            if index >= 200:
                break
    return "\n".join(rows)


def _extract_human_readable_content(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md"}:
        return _read_text_file(path)
    if suffix == ".json":
        try:
            data = json.loads(_read_text_file(path))
            return json.dumps(data, ensure_ascii=False, indent=2)
        except json.JSONDecodeError:
            return _read_text_file(path)
    if suffix == ".csv":
        return _extract_from_csv(path)
    if suffix == ".pdf":
        return _extract_from_pdf(path)
    if suffix == ".docx":
        return _extract_from_docx(path)
    if suffix == ".xlsx":
        return _extract_from_xlsx(path)
    if suffix == ".pptx":
        return _extract_from_pptx(path)
    return _read_text_file(path)


def store_uploaded_file(content: bytes, filename: str) -> Dict[str, Any]:
    suffix = _ensure_allowed_extension(filename)
    file_size = len(content)
    if file_size > settings.max_file_size:
        raise HTTPException(status_code=413, detail="文件大小超过限制")

    attachment_id = _generate_attachment_id()
    attachment_dir = ATTACHMENTS_DIR / attachment_id
    attachment_dir.mkdir(parents=True, exist_ok=True)

    sanitized_filename = _sanitize_filename(filename)
    file_path = attachment_dir / sanitized_filename
    file_path.write_bytes(content)

    readable_content = _extract_human_readable_content(file_path)
    truncated_content = readable_content[: settings.file_context_max_chars]
    (attachment_dir / "content.txt").write_text(truncated_content, encoding="utf-8")

    metadata: Dict[str, Any] = {
        "id": attachment_id,
        "filename": filename,
        "stored_filename": sanitized_filename,
        "size": file_size,
        "extension": suffix,
        "mime_type": mimetypes.guess_type(filename)[0] or "application/octet-stream",
        "preview": truncated_content[: settings.file_preview_chars],
        "uploaded_at": datetime.utcnow().isoformat() + "Z",
    }

    (attachment_dir / "meta.json").write_text(
        json.dumps(metadata, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    return metadata


def get_attachment_metadata(attachment_id: str) -> Dict[str, Any]:
    meta_path = ATTACHMENTS_DIR / attachment_id / "meta.json"
    if not meta_path.exists():
        raise HTTPException(status_code=404, detail="附件不存在或已删除")
    return json.loads(meta_path.read_text(encoding="utf-8"))


def get_attachment_content(attachment_id: str) -> str:
    content_path = ATTACHMENTS_DIR / attachment_id / "content.txt"
    if not content_path.exists():
        raise HTTPException(status_code=404, detail="附件内容不存在或已删除")
    return content_path.read_text(encoding="utf-8")


def ensure_attachments_exist(attachment_ids: List[str]) -> List[Dict[str, Any]]:
    attachments: List[Dict[str, Any]] = []
    for attachment_id in attachment_ids:
        metadata = get_attachment_metadata(attachment_id)
        content = get_attachment_content(attachment_id)
        attachments.append({**metadata, "content": content})
    return attachments